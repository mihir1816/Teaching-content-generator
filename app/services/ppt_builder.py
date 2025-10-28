"""
PowerPoint builder service with improved formatting and content extraction.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pathlib import Path
import logging

logger = logging.getLogger(__name__)


def build_ppt_from_result(result: dict, output_dir: Path = None) -> Path:
    """
    Build a PowerPoint presentation from generation result.
    
    Args:
        result: Dictionary with notes, summary, mcqs
        output_dir: Directory to save PPT (default: data/outputs)
    
    Returns:
        Path to generated PPT file
    """
    if output_dir is None:
        output_dir = Path("data/outputs")
    
    output_dir.mkdir(parents=True, exist_ok=True)
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    topic = result.get("topic", "Teaching Content")
    if isinstance(topic, list):
        topic = ", ".join(topic[:3])  # Take first 3 topics
    
    level = result.get("level", "beginner")
    
    logger.info(f"Building PPT for topic: {topic}")
    logger.info(f"Result structure: {result.keys()}")
    
    # 1. Title Slide
    _add_title_slide(prs, topic, level)
    
    # 2. Summary Slides
    summary_data = result.get("summary", {})
    if summary_data and isinstance(summary_data, dict):
        logger.info(f"Adding summary slides: {summary_data.keys()}")
        _add_summary_slides(prs, summary_data)
    
    # 3. Notes Slides
    notes_data = result.get("notes", {})
    if notes_data and isinstance(notes_data, dict):
        logger.info(f"Adding notes slides: {notes_data.keys()}")
        _add_notes_slides(prs, notes_data)
    
    # 4. Glossary Slide
    glossary = None
    if notes_data and isinstance(notes_data, dict):
        glossary = notes_data.get("glossary", [])
    
    if glossary and len(glossary) > 0:
        logger.info(f"Adding glossary with {len(glossary)} terms")
        _add_glossary_slide(prs, glossary)
    
    # 5. MCQ Slides
    mcqs_data = result.get("mcqs", {})
    if mcqs_data and isinstance(mcqs_data, dict):
        questions = mcqs_data.get("questions", [])
        if questions:
            logger.info(f"Adding {len(questions)} MCQ slides")
            _add_mcq_slides(prs, questions)
    
    # Save
    topic_safe = str(topic).replace(' ', '_')[:30]
    ppt_filename = f"{topic_safe}_presentation.pptx"
    ppt_path = output_dir / ppt_filename
    prs.save(str(ppt_path))
    
    logger.info(f"PowerPoint saved: {ppt_path}")
    return ppt_path


def _add_title_slide(prs, topic, level):
    """Add title slide with proper formatting."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(31, 78, 120)  # Dark blue
    
    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.5), Inches(9), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    
    p = title_frame.paragraphs[0]
    p.text = str(topic)
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(4.5), Inches(9), Inches(0.8)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    
    p = subtitle_frame.paragraphs[0]
    p.text = f"Level: {level.capitalize()}"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(200, 200, 200)


def _add_summary_slides(prs, summary_data):
    """Add summary slides with proper text wrapping."""
    logger.info(f"Summary data: {summary_data}")
    
    summary_text = summary_data.get("summary", "")
    key_points = summary_data.get("key_points", [])
    
    # Summary text slide
    if summary_text and len(summary_text.strip()) > 0:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(9), Inches(0.6)
        )
        title_frame = title_box.text_frame
        p = title_frame.paragraphs[0]
        p.text = "Summary"
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = RGBColor(31, 78, 120)
        
        # Content with proper bounds
        content_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(1.2), Inches(8.6), Inches(5.8)
        )
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        content_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        
        # Split long text into paragraphs
        paragraphs = summary_text.split('. ')
        for i, para in enumerate(paragraphs[:4]):  # Limit to 4 sentences per slide
            if i > 0:
                content_frame.add_paragraph()
            p = content_frame.paragraphs[i]
            para_text = para.strip()
            if para_text and not para_text.endswith('.'):
                para_text += '.'
            p.text = _truncate_text(para_text, 600)
            p.font.size = Pt(16)
            p.space_after = Pt(12)
            p.level = 0
        
        logger.info("Added summary slide")
    
    # Key points slide
    if key_points and len(key_points) > 0:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(9), Inches(0.6)
        )
        title_frame = title_box.text_frame
        p = title_frame.paragraphs[0]
        p.text = "Key Points"
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = RGBColor(31, 78, 120)
        
        # Bullet points with proper spacing
        content_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(1.2), Inches(8.6), Inches(5.8)
        )
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        for i, point in enumerate(key_points[:6]):  # Max 6 points
            if i > 0:
                content_frame.add_paragraph()
            p = content_frame.paragraphs[i]
            p.text = "• " + _truncate_text(str(point), 150)
            p.font.size = Pt(18)
            p.space_after = Pt(16)
            p.level = 0
        
        logger.info(f"Added key points slide with {len(key_points)} points")


def _add_notes_slides(prs, notes_data):
    """Add notes slides with sections - handles 'bullets' field."""
    logger.info(f"Notes data keys: {notes_data.keys()}")
    
    # Get sections
    sections = notes_data.get("sections", [])
    
    if not sections or len(sections) == 0:
        logger.warning("No sections found in notes")
        return
    
    logger.info(f"Processing {len(sections)} sections")
    
    for idx, section in enumerate(sections[:8]):  # Max 8 sections
        logger.info(f"Section {idx}: {section.get('title', 'Unknown')}")
        
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Section title
        section_title = section.get("title", f"Section {idx + 1}")
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(9), Inches(0.6)
        )
        title_frame = title_box.text_frame
        p = title_frame.paragraphs[0]
        p.text = _truncate_text(str(section_title), 80)
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(31, 78, 120)
        
        # Section content - check for 'bullets' field first, then 'content'
        bullets = section.get("bullets", [])
        
        if not bullets:
            # Fallback to 'content' field if 'bullets' doesn't exist
            content = section.get("content", "")
            if content:
                bullets = [content]
            else:
                bullets = ["No content available."]
                logger.warning(f"Empty content for section: {section_title}")
        
        # Create text box for bullets
        content_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(1.2), Inches(8.6), Inches(5.8)
        )
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        # Add each bullet point
        for i, bullet in enumerate(bullets[:6]):  # Max 6 bullets per slide
            if i > 0:
                content_frame.add_paragraph()
            p = content_frame.paragraphs[i]
            p.text = "• " + _truncate_text(str(bullet), 250)
            p.font.size = Pt(16)
            p.space_after = Pt(14)
            p.level = 0
        
        logger.info(f"Added section slide: {section_title} with {len(bullets)} bullets")


def _add_glossary_slide(prs, glossary):
    """Add glossary slide in single-column layout."""
    # Create multiple slides if glossary is too long
    terms_per_slide = 8
    total_slides = (len(glossary) + terms_per_slide - 1) // terms_per_slide
    
    for slide_num in range(total_slides):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(9), Inches(0.6)
        )
        title_frame = title_box.text_frame
        p = title_frame.paragraphs[0]
        title_text = "Glossary" if total_slides == 1 else f"Glossary ({slide_num + 1}/{total_slides})"
        p.text = title_text
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = RGBColor(31, 78, 120)
        
        # Content
        content_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(1.2), Inches(8.6), Inches(5.8)
        )
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        # Get terms for this slide
        start_idx = slide_num * terms_per_slide
        end_idx = min(start_idx + terms_per_slide, len(glossary))
        terms_for_slide = glossary[start_idx:end_idx]
        
        for i, item in enumerate(terms_for_slide):
            if i > 0:
                content_frame.add_paragraph()
            
            term = item.get("term", "")
            definition = item.get("definition", "")
            
            p = content_frame.paragraphs[i]
            p.text = f"• {term}: {_truncate_text(definition, 150)}"
            p.font.size = Pt(14)
            p.space_after = Pt(12)


def _add_mcq_slides(prs, questions):
    """Add MCQ slides - handles 'stem' and 'answer' fields."""
    for idx, q in enumerate(questions[:10], 1):  # Max 10 MCQs
        # Question slide
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Question number/title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(9), Inches(0.5)
        )
        title_frame = title_box.text_frame
        p = title_frame.paragraphs[0]
        p.text = f"Question {idx}"
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = RGBColor(31, 78, 120)
        
        # Question text - check for 'stem' first, then 'question'
        question_text = q.get("stem", q.get("question", ""))
        
        q_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(1.0), Inches(8.6), Inches(1.5)
        )
        q_frame = q_box.text_frame
        q_frame.word_wrap = True
        p = q_frame.paragraphs[0]
        p.text = _truncate_text(question_text, 250)
        p.font.size = Pt(18)
        p.font.bold = True
        
        # Options
        options_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(2.8), Inches(8.6), Inches(3.5)
        )
        options_frame = options_box.text_frame
        options_frame.word_wrap = True
        
        options = q.get("options", [])
        for i, opt in enumerate(options[:4]):  # Max 4 options
            if i > 0:
                options_frame.add_paragraph()
            p = options_frame.paragraphs[i]
            # Options might already have labels like "A)" or just be plain text
            opt_text = str(opt).strip()
            if not opt_text.startswith(('A)', 'B)', 'C)', 'D)')):
                label = chr(65 + i)  # A, B, C, D
                opt_text = f"{label}) {opt_text}"
            p.text = _truncate_text(opt_text, 150)
            p.font.size = Pt(16)
            p.space_after = Pt(12)
        
        # Answer slide
        answer_slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Answer title
        ans_title_box = answer_slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(9), Inches(0.5)
        )
        ans_title_frame = ans_title_box.text_frame
        p = ans_title_frame.paragraphs[0]
        p.text = f"Answer - Question {idx}"
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = RGBColor(46, 125, 50)
        
        # Correct answer - check for 'answer' or 'correct_answer'
        correct_ans = q.get("answer", q.get("correct_answer", ""))
        
        ans_box = answer_slide.shapes.add_textbox(
            Inches(0.7), Inches(1.5), Inches(8.6), Inches(1.0)
        )
        ans_frame = ans_box.text_frame
        ans_frame.word_wrap = True
        p = ans_frame.paragraphs[0]
        p.text = f"✓ Correct Answer: {correct_ans}"
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = RGBColor(46, 125, 50)
        
        # Explanation
        explanation = q.get("explanation", "No explanation provided.")
        exp_box = answer_slide.shapes.add_textbox(
            Inches(0.7), Inches(3.0), Inches(8.6), Inches(3.5)
        )
        exp_frame = exp_box.text_frame
        exp_frame.word_wrap = True
        exp_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        p = exp_frame.paragraphs[0]
        p.text = _truncate_text(explanation, 500)
        p.font.size = Pt(16)
        p.space_after = Pt(10)


def _truncate_text(text: str, max_chars: int) -> str:
    """Truncate text to max characters, adding ellipsis if needed."""
    if not text:
        return ""
    
    text = str(text).strip()
    if len(text) <= max_chars:
        return text
    
    # Find last space before max_chars
    truncated = text[:max_chars]
    last_space = truncated.rfind(' ')
    
    if last_space > max_chars * 0.8:  # If space is reasonably close
        return truncated[:last_space] + "..."
    
    return truncated + "..."